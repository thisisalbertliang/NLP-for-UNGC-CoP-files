#fnx and fnc are file path for final output
fnx='PLEASE INSERT YOUR DESIRED OUTPUT FOLDER HERE/2018_COP_phase_1_result.xlsx'
fnc='PLEASE INSERT YOUR DESIRED OUTPUT FOLDER HERE/2018_COP_phase_1_result.csv'

#Import all the necessary modules and packages
import time
import sys
import warnings
import re
import json

import img2pdf
from PIL import Image
import textract
import PyPDF2
import os
from langdetect import detect
import time
import numpy as np
import pandas as pd
import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)
import textwrap
import requests
import uuid
import glob
#### Time out
import signal
from time import sleep    # only needed for testing
import docx

# Custom exception for the timeout
class TimeoutException(Exception):
    pass

# Handler function to be called when SIGALRM is received
def sigalrm_handler(signum, frame):
    # We get signal!
    raise TimeoutException()

##### Convert Docx and PPTX: PPTX returns a list
from pptx import Presentation
def getPPT(path_to_presentation):
    prs = Presentation(path_to_presentation)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    slide_num = 1
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        slide_num += 1
    return [slide_num, text_runs]

def getText(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def get_translated_text(text):
    body = [{'text': text}]
    request = requests.post(constructed_url, headers=headers, json=body)
    response = request.json()
    #return response
    return response[0]["translations"][0]["text"]

def pdf2text_and_numpage(filename):
    #write a for-loop to open many files
    #open allows you to read the file
    pdfFileObj = open(filename,'rb')
    #The pdfReader variable is a readable object that will be parsed
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    #discerning the number of pages will allow us to parse through all #the pages
    num_pages = pdfReader.numPages
    count = 0
    text = ""
    #The while loop will read each page
    while count < num_pages:
        pageObj = pdfReader.getPage(count)
        count +=1
        text += pageObj.extractText()
    #This if statement exists to check if the above library returned #words. It's done because PyPDF2 cannot read scanned files.
    if text != "":
        return (text.decode('utf-8'), num_pages, 'text-based-PDF')
    #If the above returns as False, we run the OCR library textract to #convert scanned/image based PDF files into text
    else:
        return ((textract.process(filename, method='tesseract')).decode('utf-8'), num_pages, 'image-based-PDF')
    # Now we have a text variable which contains all the text derived #from our PDF file. Type print(text) to see what it contains. It #likely contains a lot of spaces, possibly junk such as '\n' etc.
    # Now, we will clean our text variable, and return it as a list of keywords.

def image2text_and_numpage(img_path, path2original):
    # storing pdf path
    pdf_path = path2original + 'converted.pdf'
    # opening image
    image = Image.open(img_path)
    # converting into chunks using img2pdf
    pdf_bytes = img2pdf.convert(image.filename)
    # opening or creating pdf file
    file = open(pdf_path, "wb")
    # writing pdf files with chunks
    file.write(pdf_bytes)
    # closing image file
    image.close()
    # closing pdf file
    file.close()
    return pdf2text_and_numpage(pdf_path)

## save all the folders name of articles into the corp_id list
directory = r'PLEASE INSERT THE PATH TO THE FOLDER WHERE YOU STORED ALL THE CORPID FOLDERS'
os.chdir(directory)
#print(os.getcwd())
corp_id = [x for x in glob.glob('*')]

#This step is to save all the name and id of articles in list1 and files_name list

#num is used to calculate the number of empty folders
num = 0
#a is used to calculate the number of non-empty folders
a = 0
#list1 is used to save all the id of non-empty articles
list1 = []
#files_name is a list used to save all the names of non-empty articles
files_name = []

for i in corp_id:
    ##each path is the path of each articles
    path = directory + '/' + i + '/original'
    for dirpath, dirnames, files in os.walk(path):
        if not files:
            #print(dirpath, 'is empty')
            num += 1
        else:
            a += 1
            #print(dirpath,'is not empty')
            list1.append(i)
            files_name.append(files)

#generate a dataframe which includes the corporate_id and size
#d=pd.DataFrame(columns=['corpor_id','file_size_in_bytes'])
#result_list is a list of list, which contains all id, file names and size of all the elements of articles
result_list = []
fail_file = []
for i in range(0, len(files_name)):
    if len(files_name[i]) == 1:
        ## file_name is the first element of each item in files_name list
        file_name = str(files_name[i][0])
        try:
            #### list1 is id
            path1 = directory + '/' + str(list1[i]) + '/original/' + file_name
            ## statinfo is the size of each file Size in bytes of a plain file; amount of data waiting on some special files.
            statinfo = os.stat(path1).st_size
            result_list.append([list1[i],file_name,statinfo])
        except:
            fail_file.append([i,list1[i],file_name])
            pass
    else:
        for l in range(0,len(files_name[i])):
            file_name = str(files_name[i][l])
            try:
                path1 = directory + '/' + str(list1[i]) + '/original/' + file_name
                statinfo = os.stat(path1).st_size
                result_list.append([list1[i],file_name,statinfo])
            except:
                fail_file.append([i,list1[i],file_name])
                pass

labels = ['corp_id', 'file_name', 'file_size_in_bytes']
#generate a dataframe called df_result with column names 'corp_id','file_name','size'
#the elements in df_result is result_list
df_result = pd.DataFrame(result_list, columns=labels)

####LIMIT FOR CLOUD JSON IS DIFFERENT
unit_limit= 10000000
#os.environ["GOOGLE_APPLICATION_CREDENTIALS"]='/home/globalaicloud/UNGC_NEW/JSON/My First Project-dec55926c2d0.json'

# Checks to see if the Translator Text subscription key is available
# as an environment variable. If you are setting your subscription key as a
# string, then comment these lines out.
if 'TRANSLATOR_TEXT_KEY' in os.environ:
    subscriptionKey = os.environ['TRANSLATOR_TEXT_KEY']
else:
    print('Environment variable for TRANSLATOR_TEXT_KEY is not set.')
    #exit()
# If you want to set your subscription key as a string, uncomment the line
# below and add your subscription key.
subscriptionKey = "331f7ace25a849639d0d319181758dff"

base_url = 'https://api.cognitive.microsofttranslator.com'
path = '/translate?api-version=3.0'
params = '&to=en'
constructed_url = base_url + path + params

headers = {
    'Ocp-Apim-Subscription-Key': subscriptionKey,
    'Content-type': 'application/json',
    'X-ClientTraceId': str(uuid.uuid4())}

timelimit_seconds = 60*30    # Must be an integer


### fill the columns in each row
###"file format" "page number" "original text" ""
start = time.time()
print('Start time: {}'.format(start))

for i in range(len(result_list)):
    path1 = directory + '/' + str(df_result.corp_id[i]) + '/original/' + df_result.file_name[i]
    path2original = directory + '/' + str(df_result.corp_id[i]) + '/original/'
    file_name = str(df_result.file_name[i])
    if file_name.endswith('.doc') or file_name.endswith('.docx'):
        df_result.loc[i, 'file_format'] = 'Word File'
        df_result.loc[i, 'page_num'] = np.NaN
        try:
            df_result.loc[i, 'text'] = getText(path1)
            print(str(i), df_result.corp_id[i], file_name, 'Doc converted')
        except:
            df_result.loc[i, 'text'] = 'Wrong Format'
            print(str(i), df_result.corp_id[i], file_name, 'Doc unconverted')
    elif file_name.endswith('.ppt') or file_name.endswith('.pptx'):
        df_result.loc[i, 'file_format'] = 'Powerpoint'
        try:
            ppt_text = getPPT(path1)
            df_result.loc[i, 'page_num'] = ppt_text[0]
            df_result.loc[i, 'text'] = ''.join(x for x in ppt_text[1])
            print(str(i), df_result.corp_id[i], file_name, 'PPT converted')
        except:
            df_result.loc[i, 'page_num'] = np.NaN
            df_result.loc[i, 'text'] = 'Wrong Format'
            print(str(i), df_result.corp_id[i], file_name, 'PPT unconverted')
    elif file_name.endswith('.jpg') or file_name.endswith('.jpeg') or file_name.endswith('.png'):
        try:
            wildcard, file_extension = os.path.splitext(path1)
            image_text, num_page, wildcard = image2text_and_numpage(path1, path2original)
            df_result.loc[i, 'page_num'] = num_page
            df_result.loc[i, 'text'] = image_text
            df_result.loc[i, 'file_format'] = file_extension[1:]
            print(str(i), df_result.corp_id[i], file_name, 'image converted')
        except Exception as e:
            print("Error Message for IMAGE READ: ")
            print(e)
    elif file_name.endswith('.xlsx') or file_name.endswith('.csv'):
        df_result.loc[i, 'file_format'] = 'Excel'
        df_result.loc[i, 'page_num'] = np.NaN
        df_result.loc[i, 'text'] = 'Wrong Format'
        print(str(i), df_result.corp_id[i], file_name, 'Wrong Format')
    else:
        raw_str = ''
        # Set up signal handler for SIGALRM, saving previous value
        old_handler = signal.signal(signal.SIGALRM, sigalrm_handler)
        # Start timer
        signal.alarm(timelimit_seconds)
        try:
            pdf_text, num_page, pdf_type = pdf2text_and_numpage(path1)
            df_result.loc[i, 'page_num'] = num_page
            df_result.loc[i, 'text'] = pdf_text
            df_result.loc[i, 'file_format'] = pdf_type
            print(str(i), df_result.corp_id[i], file_name, 'PDF converted')
        except Exception as e:
            if str(e) is "":
                df_result.loc[i, 'text'] = 'Timeout'
                df_result.loc[i, 'translation'] = 'Timeout'
                df_result.loc[i, 'page_num'] = np.NaN
                print(str(i), df_result.corp_id[i], file_name, 'took too long to convert')
            else:
                print('ERROR MESSAGE PDF:')
                print(e)
                df_result.loc[i, 'text'] = 'Convert Manually'
                df_result.loc[i, 'page_num'] = np.NaN
                df_result.loc[i, 'translation'] = 'Error'
                print(str(i), df_result.corp_id[i], file_name, 'fail to convert')
        finally:
            # Turn off timer
            signal.alarm(0)
            # Restore handler to previous value
            signal.signal(signal.SIGALRM, old_handler)
    try:
        if detect(df_result.loc[i, 'text']) == 'en':
            df_result.loc[i, 'translation'] = 'already in English'
            print(str(i), df_result.corp_id[i], file_name, 'English')
        else:
            if len(df_result.loc[i, 'text']) >= unit_limit:
                df_result.loc[i, 'translation'] = 'Too big'
                print(str(i), df_result.corp_id[i], file_name, 'too big to translate')
            else:
                try:
                    if len(df_result.loc[i, 'text']) <= 4000:
                        df_result.loc[i, 'translation'] = get_translated_text(df_result.loc[i, 'text'])
                    else:
                        lines = textwrap.wrap(df_result.loc[i, 'text'], 4000, break_long_words=False)
                        df_result.loc[i, 'translation'] = ''.join([get_translated_text(x) for x in lines])
                    print(str(i), df_result.corp_id[i], file_name, 'Translated')
                except Exception as e:
                    df_result.loc[i, 'translation'] = 'Fail to translate'
                    print(str(i), df_result.corp_id[i], file_name, 'Fail to translate')
                    print("ERROR MESSAGE translation: ")
                    print(e)
    except Exception as e:
        print('ERROR MESSAGE detect:')
        print(e)
        df_result.loc[i, 'text'] = 'Actually image'
        df_result.loc[i, 'page_num'] = np.NaN
        df_result.loc[i, 'translation'] = 'No content'
        print(str(i), df_result.corp_id[i], file_name, 'No content')
    writer = pd.ExcelWriter(fnx)
    df_result.loc[:, :].to_excel(writer, 'sheet1', header=True, index=True)
    writer.save()
    writer.close()
    df_result.loc[:, :].to_csv(fnc, index=True)
    print(str(i), df_result.corp_id[i], ' finish step', fnc)

df_result['word_count'] = df_result['text'].apply(lambda txt: len(txt.split()) if isinstance(txt, str) else -1)

end = time.time()
print('End time: {}'.format(end))
print('Total run time: {}'.format(end - start))
